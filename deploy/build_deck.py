#!/usr/bin/env python3
"""
Generate a crisp 9-slide deck for the Agent Trajectory Analytics project.
Standalone (no template needed). Dark theme matching the dashboard.

Usage:  python3 deploy/build_deck.py [output.pptx]
"""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── palette (matches the analytics dashboard) ──
BG      = RGBColor(0x0F, 0x17, 0x2A)   # deep slate
CARD    = RGBColor(0x1E, 0x29, 0x3B)   # card slate
INK     = RGBColor(0xE2, 0xE8, 0xF0)   # near-white text
MUTED   = RGBColor(0x94, 0xA3, 0xB8)   # muted grey-blue
ACCENT  = RGBColor(0x38, 0xBD, 0xF8)   # cyan accent
GREEN   = RGBColor(0x4A, 0xDE, 0x80)
AMBER   = RGBColor(0xFB, 0xBF, 0x24)
RED     = RGBColor(0xF8, 0x71, 0x71)
PURPLE  = RGBColor(0xA7, 0x8B, 0xFA)

W, H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(1, 0, 0, W, H)
    r.fill.solid(); r.fill.fore_color.rgb = BG
    r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element)
    s.shapes._spTree.insert(2, r._element)
    return s


def box(s, x, y, w, h, fill=None, line=None):
    shp = s.shapes.add_shape(1, x, y, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp=1.0):
    """runs: list of paragraphs; each paragraph is list of (txt, size, color, bold)."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = sp
        p.space_after = Pt(4)
        for (txt, size, color, bold) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.name = "Calibri"
    return tb


def chip(s, x, y, label, color):
    c = box(s, x, y, Inches(0.32), Inches(0.32), fill=color)
    return c


def footer(s, n, label):
    text(s, Inches(0.6), Inches(7.0), Inches(9), Inches(0.4),
         [[("Agent Trajectory Analytics", 10, MUTED, False)]])
    text(s, Inches(11.4), Inches(7.0), Inches(1.3), Inches(0.4),
         [[(f"{n} / 9", 10, MUTED, False)]], align=PP_ALIGN.RIGHT)
    box(s, Inches(0.6), Inches(6.95), Inches(12.1), Pt(1.2), fill=CARD)


def heading(s, kicker, title, kicker_color=ACCENT):
    text(s, Inches(0.6), Inches(0.45), Inches(12), Inches(0.4),
         [[(kicker, 13, kicker_color, True)]])
    text(s, Inches(0.6), Inches(0.85), Inches(12.1), Inches(0.9),
         [[(title, 30, INK, True)]])


# ════════════════════════════════════════════════ SLIDE 1 — Title
s = slide()
box(s, 0, Inches(2.55), W, Pt(3), fill=ACCENT)
text(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(1.0),
     [[("Agent Trajectory Analytics", 44, INK, True)]])
text(s, Inches(0.9), Inches(2.75), Inches(11.5), Inches(0.7),
     [[("Infrastructure-aware observability for non-deterministic AI agents", 20, ACCENT, False)]])
text(s, Inches(0.9), Inches(3.65), Inches(11.5), Inches(1.4),
     [[("Trace how an agent's execution paths drift, score its answer quality", 16, MUTED, False)],
      [("with an LLM-as-judge, and ", 16, MUTED, False), ("prove whether a quality drop was caused", 16, INK, True)],
      [("by GPU contention or by the application itself.", 16, INK, True)]])
text(s, Inches(0.9), Inches(6.2), Inches(11), Inches(0.5),
     [[("Agam Jain", 15, INK, True), ("    •    OpenTelemetry · Spark · Delta Lake · Ollama (real GPU)", 13, MUTED, False)]])

# ════════════════════════════════════════════════ SLIDE 2 — Problem
s = slide()
heading(s, "THE PROBLEM", "AI agents are a black box when they fail", RED)
text(s, Inches(0.6), Inches(1.9), Inches(12), Inches(0.8),
     [[("AI agents are ", 18, INK, False), ("non-deterministic", 18, RED, True),
       (" — the same request can take a different path each run, and answer quality can quietly collapse.", 18, INK, False)]])
# two contrasting cards
box(s, Inches(0.6), Inches(3.0), Inches(5.9), Inches(2.5), fill=CARD)
text(s, Inches(0.95), Inches(3.25), Inches(5.3), Inches(2.1),
     [[("😕  The question teams ask", 16, AMBER, True)],
      [("", 6, MUTED, False)],
      [("“Was the model just dumb today?”", 17, INK, True)],
      [("", 6, MUTED, False)],
      [("Shallow. Gives no root cause. Can't tell if it's the model, the prompt, or the hardware.", 14, MUTED, False)]])
box(s, Inches(6.85), Inches(3.0), Inches(5.9), Inches(2.5), fill=CARD)
text(s, Inches(7.2), Inches(3.25), Inches(5.3), Inches(2.1),
     [[("🎯  The question we answer", 16, GREEN, True)],
      [("", 6, MUTED, False)],
      [("“Did the infrastructure underneath the agent change how it behaved?”", 17, INK, True)],
      [("", 6, MUTED, False)],
      [("Separates infra-caused failures from app-caused ones — so you fix the right thing.", 14, MUTED, False)]])
text(s, Inches(0.6), Inches(5.85), Inches(12), Inches(0.8),
     [[("Nobody looks at these three together: ", 15, MUTED, False),
       ("execution path  +  answer quality  +  GPU pressure", 15, ACCENT, True),
       ("  at the same moment.", 15, MUTED, False)]])
footer(s, 2, "Problem")

# ════════════════════════════════════════════════ SLIDE 3 — What is a trajectory
s = slide()
heading(s, "KEY CONCEPT", "What is an “agent trajectory”?")
text(s, Inches(0.6), Inches(1.85), Inches(12.1), Inches(0.7),
     [[("A trajectory is the ", 17, INK, False), ("ordered sequence of steps", 17, ACCENT, True),
       (" an agent takes to answer one request — its execution path, captured as a signature.", 17, INK, False)]])
# healthy path
box(s, Inches(0.6), Inches(2.75), Inches(12.1), Inches(1.35), fill=CARD)
text(s, Inches(0.9), Inches(2.9), Inches(11.5), Inches(0.4),
     [[("✅  Healthy path  ", 15, GREEN, True), ("(calm GPU)", 13, MUTED, False),
       ("            quality 3.5 / 5", 14, GREEN, True)]])
text(s, Inches(0.9), Inches(3.35), Inches(11.6), Inches(0.7),
     [[("ENTRY → PLAN → RESEARCH → RETRIEVE → REASON → VERIFY → RESPOND", 16, INK, True)],
      [("Full research and a verification step before answering.", 13, MUTED, False)]])
# mutated path
box(s, Inches(0.6), Inches(4.3), Inches(12.1), Inches(1.35), fill=CARD)
text(s, Inches(0.9), Inches(4.45), Inches(11.5), Inches(0.4),
     [[("⚠️  Mutated path  ", 15, RED, True), ("(GPU under pressure)", 13, MUTED, False),
       ("   quality 0.6 / 5", 14, RED, True)]])
text(s, Inches(0.9), Inches(4.9), Inches(11.6), Inches(0.7),
     [[("ENTRY → PLAN → REASON → RESPOND", 16, INK, True),
       ("      −3 steps, skipped research + verify", 13, RED, False)],
      [("Under load the agent truncated its own reasoning → wrong answer.", 13, MUTED, False)]])
text(s, Inches(0.6), Inches(5.95), Inches(12), Inches(0.7),
     [[("We hash each path into a ", 14, MUTED, False), ("template signature", 14, ACCENT, True),
       (", so we can count, compare, and detect when the agent starts drifting.", 14, MUTED, False)]])
footer(s, 3, "Concept")

# ════════════════════════════════════════════════ SLIDE 4 — Solution / pipeline
s = slide()
heading(s, "THE SOLUTION", "One pipeline that joins behavior to infrastructure", GREEN)
text(s, Inches(0.6), Inches(1.8), Inches(12.1), Inches(0.55),
     [[("A travel-planner agent is fully instrumented; the data streams through Spark + Delta Lake into a live dashboard.", 15, MUTED, False)]])
stages = [
    ("Instrument", "Every agent / LLM /\nRAG / tool step is an\nOpenTelemetry span", ACCENT),
    ("Stream", "Spark Structured\nStreaming + Delta Lake\n(6 jobs, MERGE upserts)", PURPLE),
    ("Judge", "A 2nd LLM scores each\nanswer on 5 quality\ndimensions (1–5)", AMBER),
    ("Correlate", "Join quality to the GPU\nconditions present\nwhile the trace ran", GREEN),
    ("Explain", "Dashboard: drift, alerts,\nharmful mutations,\nroot cause", RED),
]
x = Inches(0.6); cw = Inches(2.28); gap = Inches(0.13)
for i, (t, d, c) in enumerate(stages):
    bx = Emu(int(x) + i * (int(cw) + int(gap)))
    box(s, bx, Inches(2.65), cw, Inches(2.2), fill=CARD)
    box(s, bx, Inches(2.65), cw, Pt(4), fill=c)
    text(s, bx, Inches(2.95), cw, Inches(0.5),
         [[(f"{i+1}. {t}", 15, c, True)]], align=PP_ALIGN.CENTER)
    text(s, Emu(int(bx)+Inches(0.12)), Inches(3.5), Emu(int(cw)-Inches(0.24)), Inches(1.3),
         [[(line, 12, INK, False)] for line in d.split("\n")], align=PP_ALIGN.CENTER)
text(s, Inches(0.6), Inches(5.35), Inches(12.1), Inches(1.1),
     [[("Correlation rule:  ", 15, ACCENT, True),
       ("a GPU/network sample attaches to a trace only if it's on the ", 14, INK, False),
       ("same node", 14, INK, True),
       (" within ", 14, INK, False), ("±10 s", 14, INK, True),
       (" of execution —", 14, INK, False)],
      [("so a quality drop is tied to the ", 14, MUTED, False),
       ("specific GPU conditions that produced it", 14, ACCENT, True),
       (", not a global average.", 14, MUTED, False)]])
footer(s, 4, "Solution")

# ════════════════════════════════════════════════ SLIDE 5 — The 3 questions
s = slide()
heading(s, "WHAT WE SET OUT TO PROVE", "Three questions, three verdicts")
qs = [
    ("Q1", "How do agent paths\nchange over time?", "Detect and quantify drift", ACCENT),
    ("Q2", "Does a quality drop\nco-occur with path change\nAND GPU contention?", "Separate infra-caused\nfrom app-caused failures", GREEN),
    ("Q3", "Which trajectory mutations\nlead to wrong answers?", "Pinpoint the failing paths\n+ their root cause", RED),
]
x = Inches(0.6); cw = Inches(3.9); gap = Inches(0.2)
for i, (q, ttl, sub, c) in enumerate(qs):
    bx = Emu(int(x) + i * (int(cw) + int(gap)))
    box(s, bx, Inches(2.2), cw, Inches(3.6), fill=CARD)
    box(s, bx, Inches(2.2), cw, Pt(5), fill=c)
    text(s, bx, Inches(2.5), cw, Inches(0.7),
         [[(q, 30, c, True)]], align=PP_ALIGN.CENTER)
    text(s, Emu(int(bx)+Inches(0.2)), Inches(3.35), Emu(int(cw)-Inches(0.4)), Inches(1.4),
         [[(l, 16, INK, True)] for l in ttl.split("\n")], align=PP_ALIGN.CENTER)
    text(s, Emu(int(bx)+Inches(0.2)), Inches(4.85), Emu(int(cw)-Inches(0.4)), Inches(0.9),
         [[(l, 13, MUTED, False)] for l in sub.split("\n")], align=PP_ALIGN.CENTER)
footer(s, 5, "Goals")

# ════════════════════════════════════════════════ SLIDE 6 — Q1 proved
s = slide()
heading(s, "PROVED · Q1", "Paths drift — and we measure exactly how much", ACCENT)
text(s, Inches(0.6), Inches(1.85), Inches(12.1), Inches(0.9),
     [[("Method: ", 16, ACCENT, True),
       ("bucket traces into time windows and compute ", 15, INK, False),
       ("Jensen-Shannon divergence", 15, INK, True),
       (" of each window's path mix vs. a calm baseline.", 15, INK, False)]])
# drift bar progression
labels = ["W1", "W2", "W3", "W4", "W5"]
vals   = [0.00, 0.13, 0.47, 0.59, 1.00]
bx0 = Inches(1.0); bw = Inches(1.7); gp = Inches(0.55); base = Inches(5.4); maxh = Inches(2.4)
for i, (lb, v) in enumerate(zip(labels, vals)):
    bxx = Emu(int(bx0) + i * (int(bw) + int(gp)))
    hh = Emu(int(Inches(0.05)) + int(maxh * v))
    yy = Emu(int(base) - int(hh))
    col = GREEN if v < 0.3 else (AMBER if v < 0.6 else RED)
    box(s, bxx, yy, bw, hh, fill=col)
    text(s, bxx, Emu(int(yy)-Inches(0.42)), bw, Inches(0.4),
         [[(f"{v:.2f}", 15, INK, True)]], align=PP_ALIGN.CENTER)
    text(s, bxx, Emu(int(base)+Inches(0.05)), bw, Inches(0.35),
         [[(lb, 13, MUTED, False)]], align=PP_ALIGN.CENTER)
text(s, Inches(0.6), Inches(6.1), Inches(12), Inches(0.6),
     [[("Result:  drift climbs ", 15, INK, False), ("0.00 → 1.00", 15, ACCENT, True),
       (" across windows — one dominant path in the calm phase fragments into many under GPU stress.", 15, INK, False)]])
footer(s, 6, "Q1 Drift")

# ════════════════════════════════════════════════ SLIDE 7 — Q2 proved (centerpiece)
s = slide()
heading(s, "PROVED · Q2  ★ CENTERPIECE", "We tell infra-caused from app-caused failures", GREEN)
text(s, Inches(0.6), Inches(1.85), Inches(12.1), Inches(0.65),
     [[("Method: ", 16, GREEN, True),
       ("per window, check three independent signals — ", 15, INK, False),
       ("quality drop · path drift · GPU pressure", 15, INK, True), (".", 15, INK, False)]])
# verdict card A
box(s, Inches(0.6), Inches(2.7), Inches(5.95), Inches(2.7), fill=CARD)
box(s, Inches(0.6), Inches(2.7), Pt(6), Inches(2.7), fill=RED)
text(s, Inches(0.95), Inches(2.95), Inches(5.4), Inches(2.3),
     [[("gpu_induced_degradation", 17, RED, True)],
      [("", 5, MUTED, False)],
      [("All three signals fired together.", 14, INK, False)],
      [("quality −1.9   ·   drift high   ·   GPU 0.71", 14, MUTED, False)],
      [("", 5, MUTED, False)],
      [("→ The infrastructure caused it.", 15, INK, True)]])
# verdict card B
box(s, Inches(6.85), Inches(2.7), Inches(5.95), Inches(2.7), fill=CARD)
box(s, Inches(6.85), Inches(2.7), Pt(6), Inches(2.7), fill=AMBER)
text(s, Inches(7.2), Inches(2.95), Inches(5.4), Inches(2.3),
     [[("app_layer_degradation", 17, AMBER, True)],
      [("", 5, MUTED, False)],
      [("Quality crashed but GPU was idle.", 14, INK, False)],
      [("quality −1.96   ·   drift yes   ·   GPU 0.21", 14, MUTED, False)],
      [("", 5, MUTED, False)],
      [("→ A prompt / model issue, NOT infra.", 15, INK, True)]])
text(s, Inches(0.6), Inches(5.7), Inches(12.1), Inches(1.0),
     [[("Why it matters:  ", 15, GREEN, True),
       ("a dumb monitor pages the infra team for both. Ours says ", 14, INK, False),
       ("which", 14, INK, True),
       (" failures are actually hardware — so the right team fixes the right thing.", 14, INK, False)]])
footer(s, 7, "Q2 Correlation")

# ════════════════════════════════════════════════ SLIDE 8 — Q3 proved
s = slide()
heading(s, "PROVED · Q3", "We pinpoint the harmful mutations — and why", RED)
text(s, Inches(0.6), Inches(1.85), Inches(12.1), Inches(0.65),
     [[("Method: ", 16, RED, True),
       ("group traces by path signature; compute bad-rate, lift vs. baseline, and step/LLM deltas vs. the healthy path.", 15, INK, False)]])
# mini table
rows = [
    ("Path", "Bad rate", "Lift", "Δ steps", "Δ LLM", "Quality"),
    ("Healthy dominant", "0%", "—", "0", "0", "3.56"),
    ("Mutation A", "100%", "3.4×", "−6", "−1", "~0"),
    ("Mutation B", "100%", "3.4×", "−6", "−1", "~0"),
]
tx = Inches(0.6); ty = Inches(2.75); tw = Inches(12.1); rh = Inches(0.62)
colw = [Inches(3.7), Inches(1.7), Inches(1.5), Inches(1.7), Inches(1.7), Inches(1.8)]
for ri, row in enumerate(rows):
    cy = Emu(int(ty) + ri * int(rh))
    fill = CARD if ri == 0 else (BG if ri % 2 else RGBColor(0x16, 0x20, 0x30))
    cx = int(tx)
    for ci, val in enumerate(row):
        box(s, Emu(cx), cy, colw[ci], rh, fill=(CARD if ri == 0 else fill), line=BG)
        is_head = ri == 0
        col = MUTED if is_head else INK
        if ri >= 2 and ci in (1, 2):
            col = RED
        if ri == 1 and ci in (1, 5):
            col = GREEN
        text(s, Emu(cx+int(Inches(0.15))), Emu(int(cy)+int(Inches(0.12))), colw[ci], rh,
             [[(val, 13, col, is_head or ci == 0)]])
        cx += int(colw[ci])
text(s, Inches(0.6), Inches(5.5), Inches(12.1), Inches(1.2),
     [[("Result:  ", 15, RED, True),
       ("the failing paths run at ", 14, INK, False), ("100% bad rate (3.4× baseline)", 14, RED, True),
       (", taking ", 14, INK, False), ("6 fewer steps and 1 fewer LLM call", 14, INK, True),
       (" than the healthy path.", 14, INK, False)],
      [("Under GPU pressure the agent truncated its reasoning — a ", 14, MUTED, False),
       ("root cause engineering can fix", 14, ACCENT, True), (", not just an alarm.", 14, MUTED, False)]])
footer(s, 8, "Q3 Mutations")

# ════════════════════════════════════════════════ SLIDE 9 — Summary
s = slide()
box(s, 0, Inches(2.4), W, Pt(3), fill=ACCENT)
text(s, Inches(0.9), Inches(0.7), Inches(11.5), Inches(0.9),
     [[("One pipeline, three answers", 34, INK, True)]])
pts = [
    ("Q1 — how", "agent paths drift over time", "drift score 0.00 → 1.00", ACCENT),
    ("Q2 — prove", "a quality drop is infra-caused vs. app-caused", "two distinct verdicts", GREEN),
    ("Q3 — pinpoint", "the exact mutations that cause wrong answers", "100% bad, 3.4× lift, −6 steps", RED),
]
y = Inches(2.7)
for (k, d, m, c) in pts:
    chip(s, Inches(0.9), y, k, c)
    text(s, Inches(1.45), Emu(int(y)-Inches(0.02)), Inches(11), Inches(0.5),
         [[(k + "  ", 16, c, True), (d, 16, INK, False), ("   —  " + m, 14, MUTED, False)]])
    y = Emu(int(y) + int(Inches(0.62)))
text(s, Inches(0.9), Inches(5.0), Inches(11.6), Inches(1.0),
     [[("From ", 17, MUTED, False), ("“the agent feels unreliable”", 17, AMBER, True),
       (" to ", 17, MUTED, False),
       ("“this truncated path, under this measured GPU contention, is exactly why — and here's the fix.”", 17, INK, True)]])
text(s, Inches(0.9), Inches(6.2), Inches(11.6), Inches(0.6),
     [[("Real, not simulated:  ", 13, GREEN, True),
       ("for the demo the full analytics pipeline ran locally; only inference + the judge ran on a real RTX 3060 (NVML metrics).", 13, MUTED, False)]])

out = sys.argv[1] if len(sys.argv) > 1 else "Agent_Trajectory_Analytics.pptx"
prs.save(out)
print(f"Saved {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
