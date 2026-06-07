#!/usr/bin/env python3
"""Trajectory Analytics Workflow — Generated presentation."""

import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.oxml.ns import nsdecls
from pptx.oxml import parse_xml
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

QN_RID = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'

# ── Allowed color values for cell_fill_xml (hex RGB only) ──
BRAND_COLORS = {
    "green": "01A982",
    "dark_slate": "444444",
    "red_orange": "CC3300",
    "white": "FFFFFF",
    "light_gray": "F2F2F2",
}


def resolve_template(explicit_path=None):
    """Find the PPTX template using explicit path or env var, or abort."""
    check_path = explicit_path or (sys.argv[1] if len(sys.argv) > 1 else None)
    if check_path:
        path = os.path.expanduser(check_path)
        if os.path.isfile(path):
            return path
        print(f"ERROR: Template not found: {path}", file=sys.stderr)
        sys.exit(1)

    env_path = os.environ.get("PPTX_TEMPLATE")
    if env_path:
        path = os.path.expanduser(env_path)
        if os.path.isfile(path):
            return path
        print(f"ERROR: PPTX_TEMPLATE is set but file not found: {path}", file=sys.stderr)
        sys.exit(1)

    print("ERROR: No PPTX template provided.\n"
          "  Option 1: python3 generate_presentation.py /path/to/template.pptx\n"
          "  Option 2: export PPTX_TEMPLATE=/path/to/template.pptx",
          file=sys.stderr)
    sys.exit(1)


def discover_layouts(prs):
    """Print every layout and its placeholders across all slide masters."""
    result = {}
    for master_idx, slide_master in enumerate(prs.slide_masters):
        for layout_idx, layout in enumerate(slide_master.slide_layouts):
            layout_key = f"{master_idx}:{layout_idx}:{layout.name}"
            print(f"\n  Master {master_idx}, Layout {layout_idx}: '{layout.name}'")
            ph_map = {}
            for ph in layout.placeholders:
                pht = ph.placeholder_format.type
                type_name = str(pht).split(".")[-1] if pht else "UNKNOWN"
                w = f"{ph.width / 914400:.1f}in" if ph.width else "?"
                h = f"{ph.height / 914400:.1f}in" if ph.height else "?"
                print(f"    PH {ph.placeholder_format.idx:>2d}  "
                      f"type={type_name:<16s}  size={w} × {h}  "
                      f"name='{ph.name}'")
                ph_map[ph.placeholder_format.idx] = (type_name, ph.width, ph.height)
            result[layout_key] = ph_map
    return result


def cell_fill_xml(color_name):
    """Return a solidFill XML element for a named brand color."""
    hex_val = BRAND_COLORS.get(color_name)
    if hex_val is None:
        raise ValueError(f"Unknown color '{color_name}'. Allowed: {list(BRAND_COLORS.keys())}")
    return parse_xml(f'<a:solidFill {nsdecls("a")}><a:srgbClr val="{hex_val}"/></a:solidFill>')


def delete_all_slides(prs):
    """Remove every slide from the presentation, keeping masters + layouts."""
    pres_elem = prs.part._element
    nsmap = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
    sldIdLst = pres_elem.find('.//p:sldIdLst', nsmap)
    if sldIdLst is None:
        return
    for sldId in list(sldIdLst):
        rId = sldId.get(QN_RID)
        sldIdLst.remove(sldId)
        if rId:
            prs.part.drop_rel(rId)


def get_layout(prs, name):
    """Find a slide layout by name or {master_idx}:{layout_idx} selector."""
    available = []

    if isinstance(name, str):
        parts = name.split(":")
        if len(parts) == 2 and all(part.isdigit() for part in parts):
            master_idx, layout_idx = (int(part) for part in parts)
            try:
                return prs.slide_masters[master_idx].slide_layouts[layout_idx]
            except IndexError as exc:
                raise KeyError(f"Layout selector '{name}' not found.") from exc

    for master_idx, master in enumerate(prs.slide_masters):
        for layout_idx, layout in enumerate(master.slide_layouts):
            available.append(f"{master_idx}:{layout_idx}={layout.name}")
            if layout.name == name:
                return layout

    raise KeyError(f"Layout '{name}' not found. Available: {available}")


def add_bullet_paragraphs(text_frame, items):
    """Add bullet-point paragraphs to a text frame."""
    for i, item in enumerate(items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.level = 0
        if isinstance(item, tuple):
            bold_text, normal_text = item
            run1 = p.add_run()
            run1.text = bold_text
            run1.font.bold = True
            if normal_text:
                run2 = p.add_run()
                run2.text = normal_text
        else:
            run = p.add_run()
            run.text = item


def add_sub_bullet(text_frame, text, level=1, bold_prefix=None):
    """Add a sub-bullet paragraph."""
    p = text_frame.add_paragraph()
    p.level = level
    if bold_prefix:
        r1 = p.add_run()
        r1.text = bold_prefix
        r1.font.bold = True
        r2 = p.add_run()
        r2.text = text
    else:
        run = p.add_run()
        run.text = text
    return p


def set_footer_and_number(slide, footer_text, slide_num):
    """Set footer placeholder and slide number if present."""
    for ph in slide.placeholders:
        if ph.placeholder_format.type.name == 'FOOTER':
            ph.text = footer_text
        elif ph.placeholder_format.type.name == 'SLIDE_NUMBER':
            ph.text = str(slide_num)


def resolve_output_path(explicit_path=None):
    """Resolve output path with security checks."""
    REJECTED_PREFIXES = ("/tmp", "/var/tmp", "/private/tmp")
    script_dir = os.path.realpath(os.path.dirname(os.path.abspath(__file__)))
    home_dir = os.path.realpath(os.path.abspath(os.path.expanduser("~")))
    workspace_root = os.environ.get("GITHUB_WORKSPACE") or os.environ.get("WORKSPACE_DIR")
    resolved_workspace_root = (
        os.path.realpath(os.path.abspath(workspace_root)) if workspace_root else None
    )

    allowed_roots = [home_dir]
    if resolved_workspace_root:
        allowed_roots.append(resolved_workspace_root)

    rejected_roots = tuple(
        os.path.realpath(os.path.abspath(prefix)) for prefix in REJECTED_PREFIXES
    )

    def _is_under(path, roots):
        resolved_path = os.path.realpath(path)
        for root in roots:
            resolved_root = os.path.realpath(root)
            try:
                if os.path.commonpath([resolved_path, resolved_root]) == resolved_root:
                    return True
            except ValueError:
                continue
        return False

    if _is_under(script_dir, allowed_roots):
        allowed_roots.append(script_dir)

    candidate = explicit_path or os.environ.get("PPTX_OUTPUT_PATH")
    if not candidate:
        if _is_under(script_dir, allowed_roots):
            candidate = os.path.join(script_dir, "trajectory_analytics_workflow.pptx")
        else:
            raise ValueError("No output path was provided.")

    output_path = os.path.realpath(os.path.abspath(os.path.expanduser(candidate)))

    for prefix in rejected_roots:
        if output_path.startswith(prefix + os.sep) or output_path == prefix:
            raise ValueError(f"Output path must not be under {prefix}: {output_path}")

    if not _is_under(output_path, allowed_roots):
        raise ValueError(f"Output path must resolve inside workspace or home: {output_path}")

    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    return output_path


# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════

template_path = resolve_template()
prs = Presentation(template_path)

print("=== Discovering layouts ===")
layout_map = discover_layouts(prs)

delete_all_slides(prs)

# ── Layout references (adjust selectors based on your template) ──
# Use discover_layouts output to find the right selectors for your template
try:
    cover_layout = get_layout(prs, 'Title Slide')
except KeyError:
    cover_layout = get_layout(prs, '0:0')

try:
    content_layout = get_layout(prs, 'Title and Content')
except KeyError:
    content_layout = get_layout(prs, '0:1')

try:
    two_col_layout = get_layout(prs, 'Two Content')
except KeyError:
    two_col_layout = content_layout

try:
    blank_layout = get_layout(prs, 'Blank')
except KeyError:
    blank_layout = get_layout(prs, '0:6')


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1: Title Slide
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(cover_layout)
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 0:  # Title
        ph.text = "Trajectory Analytics Workflow"
    elif ph.placeholder_format.idx == 1:  # Subtitle
        ph.text = "End-to-End Observability for Multi-Agent AI Systems\nAgam Jain"


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2: Agenda
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(content_layout)
slide.placeholders[0].text = "Agenda"
content = slide.placeholders[1]
tf = content.text_frame
add_bullet_paragraphs(tf, [
    "Objective",
    "The Gap Today",
    "What is an Agent Trajectory?",
    "Sample Analytics Queries",
    "Architecture",
    "Trajectory Analytics Workflow",
    "Demo",
])


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3: Objective
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(content_layout)
slide.placeholders[0].text = "Objective"
content = slide.placeholders[1]
tf = content.text_frame
add_bullet_paragraphs(tf, [
    ("Build full-stack observability ", "for multi-agent AI systems"),
    ("Trace ", "every LLM call, RAG retrieval, and tool invocation"),
    ("Score ", "response quality using LLM-as-judge"),
    ("Correlate ", "quality drops with trajectory patterns & infrastructure"),
])


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4: The Gap Today
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(content_layout)
slide.placeholders[0].text = "The Gap Today"
content = slide.placeholders[1]
tf = content.text_frame
add_bullet_paragraphs(tf, [
    ("No execution visibility ", "— agent call graphs are black boxes"),
    ("No quality attribution ", "— can't link bad outputs to specific paths"),
    ("No infra correlation ", "— GPU contention? Network latency? Unknown"),
])
add_sub_bullet(tf, "When something fails, you don't know WHY", level=1)
add_sub_bullet(tf, "Was it the agent logic, LLM reasoning, or infrastructure?", level=1)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5: What is an Agent Trajectory?
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(content_layout)
slide.placeholders[0].text = "What is an Agent Trajectory?"
content = slide.placeholders[1]
tf = content.text_frame
add_bullet_paragraphs(tf, [
    ("Definition: ", "The sequence of steps an agent takes to fulfill a request"),
])
add_sub_bullet(tf, "Orchestrator → Research Agent → Flight Agent → Hotel Agent → Itinerary Agent", level=1)
add_sub_bullet(tf, "Each step: LLM calls, RAG retrievals, tool invocations", level=1)

p = tf.add_paragraph()
p.level = 0
r = p.add_run()
r.text = "Trajectory Signature: "
r.font.bold = True
r2 = p.add_run()
r2.text = "hash of step sequence for pattern matching"

add_sub_bullet(tf, "Identify which paths correlate with high/low quality", level=1)
add_sub_bullet(tf, "Detect anomalous execution patterns", level=1)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6: Sample Analytics Queries
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(content_layout)
slide.placeholders[0].text = "Sample Analytics Queries"
content = slide.placeholders[1]
tf = content.text_frame
add_bullet_paragraphs(tf, [
    "\"Which trajectory patterns have lowest quality scores?\"",
    "\"What % of low-quality responses had GPU utilization > 90%?\"",
    "\"Show traces where network latency > 500ms correlated with failures\"",
    "\"Compare quality: RAG-heavy vs tool-heavy trajectories\"",
    "\"Which agent causes the most retries?\"",
])


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7: Architecture Overview
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(content_layout)
slide.placeholders[0].text = "Architecture Overview"
content = slide.placeholders[1]
tf = content.text_frame
add_bullet_paragraphs(tf, [
    ("Layer 1: ", "AI Agent Application — Multi-agent travel planner"),
])
add_sub_bullet(tf, "Orchestrator → Research, Flight, Hotel, Itinerary agents", level=1)
add_sub_bullet(tf, "ChromaDB (RAG) + Ollama (LLM) + MCP Tools", level=1)

p = tf.add_paragraph()
p.level = 0
r = p.add_run()
r.text = "Layer 2: "
r.font.bold = True
r2 = p.add_run()
r2.text = "Telemetry Collection — OpenTelemetry instrumentation"
add_sub_bullet(tf, "OTLP → Jaeger/Prometheus + Pulsar → Delta Lake", level=1)

p = tf.add_paragraph()
p.level = 0
r = p.add_run()
r.text = "Layer 3: "
r.font.bold = True
r2 = p.add_run()
r2.text = "Spark Streaming — 5 parallel processing jobs"

p = tf.add_paragraph()
p.level = 0
r = p.add_run()
r.text = "Layer 4: "
r.font.bold = True
r2 = p.add_run()
r2.text = "Analytics API + Dashboard"


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8: Spark Streaming Pipeline
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(content_layout)
slide.placeholders[0].text = "Trajectory Analytics Workflow"
content = slide.placeholders[1]
tf = content.text_frame
add_bullet_paragraphs(tf, [
    ("Stream 1: ", "Raw spans → agent_steps (classify & flatten)"),
    ("Stream 2: ", "agent_steps → trajectory_templates (signature extraction)"),
    ("Stream 3: ", "agent_steps → quality_scores (LLM-as-judge)"),
    ("Stream 4: ", "agent_steps → gpu_metrics, network_metrics, routing"),
    ("Stream 5: ", "Join all → trace_correlated (unified view)"),
])
add_sub_bullet(tf, "Delta MERGE upserts for exactly-once semantics", level=1)
add_sub_bullet(tf, "Real-time correlation of quality + trajectory + infra", level=1)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9: Key Capabilities
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(content_layout)
slide.placeholders[0].text = "Key Capabilities"
content = slide.placeholders[1]
tf = content.text_frame
add_bullet_paragraphs(tf, [
    ("OpenTelemetry ", "— traces, metrics, logs with session correlation"),
    ("LLM-as-Judge ", "— automated quality scoring (5 dimensions)"),
    ("Trajectory Templates ", "— pattern extraction & signature hashing"),
    ("Infrastructure Metrics ", "— GPU contention, network latency, routing"),
    ("Correlated Analytics ", "— single view joining all signals"),
    ("Natural Language Summaries ", "— LLM-powered insight generation"),
])


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10: Demo / Thank You
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(content_layout)
slide.placeholders[0].text = "Demo"
content = slide.placeholders[1]
tf = content.text_frame
add_bullet_paragraphs(tf, [
    ("Chat UI ", "— Multi-turn travel planning conversation"),
    ("Jaeger ", "— Distributed trace visualization"),
    ("Analytics Dashboard ", "— Quality scores, trajectories, correlations"),
    ("Topology View ", "— GPU cluster & request routing"),
])

p = tf.add_paragraph()
p.level = 0
p.space_before = Pt(20)
r = p.add_run()
r.text = "github.com/agam/trajectory-analytics-workflow"
r.font.italic = True


# ══════════════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════════════
output_path = resolve_output_path()
prs.save(output_path)

print(f"\n=== Presentation saved ===")
print(f"Output: {output_path}")
print(f"Slides: {len(prs.slides)}")
for i, s in enumerate(prs.slides):
    title = ""
    for ph in s.placeholders:
        if ph.placeholder_format.idx == 0 and ph.has_text_frame:
            title = ph.text[:50]
            break
    print(f"  Slide {i+1}: {title}")
